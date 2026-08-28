"""EV 에너지 소비량 예측 프로젝트 보고서 PPT 생성 스크립트."""

from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from joblib import load
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.model_selection import train_test_split

from energy_consumption.paths import DATA_DIR, MODEL_DIR, PROJECT_ROOT


OUT_DIR = PROJECT_ROOT / "ppt"
ASSET_DIR = OUT_DIR / "assets"
OUTPUT = OUT_DIR / "EV_에너지_소비량_예측_프로젝트_보고서.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "NanumSquare"
FONT_BOLD = "NanumSquare Bold"

INK = "15201D"
DARK = "103F34"
GREEN = "0D7256"
LIME = "B6D86F"
CREAM = "F4F2EA"
WHITE = "FFFEF9"
MUTED = "68756F"
LINE = "D9DED7"
PALE = "EDF2E7"
GOLD = "D6A234"
RED = "B34D3B"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.02,
             font=FONT, line_spacing=1.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, size=18, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.02)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    for text, color, bold in runs:
        run = p.add_run()
        run.text = text
        run.font.name = FONT_BOLD if bold else FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_slide_base(prs, section, number, title, subtitle=None, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(DARK if dark else CREAM)
    if dark:
        add_text(slide, section.upper(), 0.65, 0.35, 5, 0.25, 9, LIME, True)
        add_text(slide, f"{number:02d}", 12.15, 0.35, 0.5, 0.25, 9, LIME, True, PP_ALIGN.RIGHT)
        add_text(slide, title, 0.65, 0.78, 11.9, 0.75, 27, WHITE, True)
        if subtitle:
            add_text(slide, subtitle, 0.67, 1.5, 11.6, 0.5, 12, "B9CBC5")
    else:
        add_text(slide, section.upper(), 0.65, 0.35, 5, 0.25, 9, GREEN, True)
        add_text(slide, f"{number:02d}", 12.15, 0.35, 0.5, 0.25, 9, GREEN, True, PP_ALIGN.RIGHT)
        add_text(slide, title, 0.65, 0.78, 11.9, 0.65, 25, INK, True)
        if subtitle:
            add_text(slide, subtitle, 0.67, 1.42, 11.6, 0.42, 11, MUTED)
        add_line(slide, 0.65, 1.92, 12.68, 1.92, LINE, 1)
    add_text(slide, "EV ENERGY CONSUMPTION · ML PROJECT", 0.65, 7.16, 5, 0.18, 7, "819089" if not dark else "86A69B", True)
    return slide


def add_bullets(slide, items, x, y, w, h, size=15, color=INK, bullet_color=GREEN):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = f"●  {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(10)
        p.line_spacing = 1.15
        if p.runs:
            p.runs[0].font.color.rgb = rgb(bullet_color if idx == 0 else color)
    return box


def add_metric(slide, x, y, w, h, label, value, unit="", accent=False, note=None, inline_unit=False):
    add_rect(slide, x, y, w, h, DARK if accent else WHITE, DARK if accent else LINE, True)
    add_text(slide, label, x + 0.22, y + 0.2, w - 0.44, 0.28, 10, "B7CBC4" if accent else MUTED)
    if inline_unit and unit:
        box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.55), Inches(w - 0.4), Inches(0.62)
        )
        frame = box.text_frame
        frame.clear()
        frame.margin_left = frame.margin_right = Inches(0.02)
        frame.margin_top = frame.margin_bottom = Inches(0.02)
        p = frame.paragraphs[0]
        value_run = p.add_run()
        value_run.text = value
        value_run.font.name = FONT_BOLD
        value_run.font.size = Pt(29)
        value_run.font.bold = True
        value_run.font.color.rgb = rgb(WHITE if accent else INK)
        unit_run = p.add_run()
        unit_run.text = unit
        unit_run.font.name = FONT_BOLD
        unit_run.font.size = Pt(14)
        unit_run.font.bold = True
        unit_run.font.color.rgb = rgb(LIME if accent else MUTED)
    else:
        add_text(slide, value, x + 0.2, y + 0.55, w - 0.4, 0.62, 29, WHITE if accent else INK, True)
    if unit and not inline_unit:
        add_text(slide, unit, x + 0.22, y + 1.15, w - 0.44, 0.24, 9, "C9D8D3" if accent else MUTED)
    if note:
        add_text(slide, note, x + 0.22, y + h - 0.38, w - 0.44, 0.25, 8, LIME if accent else GREEN)


def add_image_fit(slide, path, x, y, w, h, border=True):
    path = Path(path)
    with Image.open(path) as image:
        iw, ih = image.size
    target_ratio = w / h
    image_ratio = iw / ih
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if image_ratio > target_ratio:
        crop = (1 - target_ratio / image_ratio) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    elif image_ratio < target_ratio:
        crop = (1 - image_ratio / target_ratio) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    if border:
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        frame.fill.background()
        frame.line.color.rgb = rgb(LINE)
    return pic


def add_image_contain(slide, path, x, y, w, h, border=True):
    """이미지를 자르지 않고 지정 영역 안에 전체가 보이도록 배치한다."""
    path = Path(path)
    with Image.open(path) as image:
        iw, ih = image.size
    image_ratio = iw / ih
    target_ratio = w / h
    if border:
        add_rect(slide, x, y, w, h, fill=WHITE, line=LINE)
    if image_ratio > target_ratio:
        shown_w = w
        shown_h = w / image_ratio
        shown_x = x
        shown_y = y + (h - shown_h) / 2
    else:
        shown_h = h
        shown_w = h * image_ratio
        shown_x = x + (w - shown_w) / 2
        shown_y = y
    return slide.shapes.add_picture(
        str(path), Inches(shown_x), Inches(shown_y),
        width=Inches(shown_w), height=Inches(shown_h)
    )


def setup_charts(df: pd.DataFrame, model_data: dict) -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    plt.rcParams.update({
        "font.family": "NanumGothic",
        "axes.unicode_minus": False,
        "figure.facecolor": WHITE,
        "axes.facecolor": WHITE,
        "text.color": INK,
        "axes.labelcolor": MUTED,
        "xtick.color": MUTED,
        "ytick.color": MUTED,
    })

    target = model_data["target_column"]
    fig, ax = plt.subplots(figsize=(8.2, 4.5))
    ax.hist(df[target], bins=34, color="#0D7256", alpha=0.9, edgecolor="#F4F2EA")
    ax.axvline(df[target].mean(), color="#B6D86F", linewidth=2.5, label=f"평균 {df[target].mean():.1f}")
    ax.set_title("에너지 소비량은 약 12~35 kWh/100km에 분포", loc="left", fontsize=16, fontweight="bold")
    ax.set_xlabel("에너지 소비량 (kWh/100km)")
    ax.set_ylabel("관측 건수")
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "target_distribution.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    corr = df.corr(numeric_only=True)[target].drop(target).sort_values()
    fig, ax = plt.subplots(figsize=(8.2, 4.9))
    colors = ["#8AB4A5" if value < 0 else "#0D7256" for value in corr]
    ax.barh(corr.index, corr.values, color=colors)
    ax.axvline(0, color="#AAB4AF", linewidth=1)
    ax.set_title("경사도·적재량·속도가 소비량과 높은 상관", loc="left", fontsize=16, fontweight="bold")
    ax.set_xlabel("Pearson 상관계수")
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.grid(axis="x", color="#E4E7E2", linewidth=.8)
    for y, value in enumerate(corr.values):
        ax.text(value + (0.012 if value >= 0 else -0.012), y, f"{value:.2f}", va="center",
                ha="left" if value >= 0 else "right", fontsize=9)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "correlations.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    model_names = ["Linear", "KNN", "KNN\nTuned", "Random\nForest", "RF\nTuned", "XGBoost", "XGBoost\nTuned"]
    rmse = [1.012057, 1.351493, 1.272699, 1.233100, 1.222125, 1.109229, 0.887273]
    fig, ax = plt.subplots(figsize=(9.2, 4.2))
    colors = ["#B8C3BE"] * 6 + ["#0D7256"]
    bars = ax.bar(model_names, rmse, color=colors, width=.68)
    ax.set_title("튜닝 XGBoost가 교차검증 RMSE 최저", loc="left", fontsize=16, fontweight="bold")
    ax.set_ylabel("CV RMSE (낮을수록 우수)")
    ax.set_ylim(0, 1.55)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.grid(axis="y", color="#E4E7E2")
    ax.set_axisbelow(True)
    for bar, value in zip(bars, rmse, strict=True):
        ax.text(bar.get_x() + bar.get_width()/2, value + .035, f"{value:.3f}", ha="center", fontsize=9, fontweight="bold")
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "model_comparison.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    features = model_data["feature_columns"]
    estimator = model_data["model"].named_steps["model"]
    importance = pd.Series(estimator.feature_importances_, index=features).sort_values()
    labels = {
        "speed_kmh": "평균 속도", "payload_kg": "적재량", "ambient_temp_C": "외기 온도",
        "hvac_power_kw": "공조 전력", "road_grade_pct": "도로 경사도", "battery_temp_C": "배터리 온도",
        "driving_style_index": "운전 성향", "tire_pressure_bar": "타이어 공기압", "trip_distance_km": "주행 거리",
    }
    fig, ax = plt.subplots(figsize=(8.3, 4.8))
    ax.barh([labels[name] for name in importance.index], importance.values * 100,
            color=["#8FA79E"] * 4 + ["#0D7256"] * 5)
    ax.set_title("도로 경사도·적재량·속도가 상위 예측 요인", loc="left", fontsize=16, fontweight="bold")
    ax.set_xlabel("Feature Importance (%)")
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.grid(axis="x", color="#E4E7E2")
    ax.set_axisbelow(True)
    for y, value in enumerate(importance.values * 100):
        ax.text(value + .3, y, f"{value:.1f}%", va="center", fontsize=9)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "feature_importance.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    X = df[features]
    y = df[target]
    _, X_test, _, y_test = train_test_split(X, y, test_size=.2, random_state=0)
    prediction = model_data["model"].predict(X_test)
    fig, ax = plt.subplots(figsize=(6.5, 5.2))
    ax.scatter(y_test, prediction, alpha=.3, s=18, color="#0D7256", edgecolors="none")
    lower = min(y_test.min(), prediction.min())
    upper = max(y_test.max(), prediction.max())
    ax.plot([lower, upper], [lower, upper], "--", color="#D6A234", linewidth=2)
    ax.set_title("실제값과 예측값이 기준선 주변에 밀집", loc="left", fontsize=16, fontweight="bold")
    ax.set_xlabel("실제 소비량 (kWh/100km)")
    ax.set_ylabel("예측 소비량 (kWh/100km)")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(color="#E8EAE6", linewidth=.7)
    fig.tight_layout()
    fig.savefig(ASSET_DIR / "actual_vs_predicted.png", dpi=180, bbox_inches="tight")
    plt.close(fig)


def build() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    df = pd.read_csv(DATA_DIR / "ev_energy_consumption.csv")
    model_data = load(MODEL_DIR / "energy_consumption_best_model.joblib")
    setup_charts(df, model_data)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 1. 표지
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(DARK)
    add_rect(slide, 8.45, 0, 4.88, 7.5, GREEN, GREEN)
    for i, width in enumerate([2.9, 3.6, 2.2, 4.0, 3.0]):
        add_rect(slide, 8.8, 1.05 + i*.9, width, .08, LIME if i in (0, 3) else "70A491", "70A491")
    add_text(slide, "ML ANALYSIS REPORT", .72, .65, 5, .3, 10, LIME, True)
    add_text(slide, "전기차 에너지 소비량\n예측 웹서비스", .72, 1.35, 7.2, 1.65, 36, WHITE, True)
    add_text(slide, "운행 조건 기반 소비량 예측과\n현재·대안 계획 비교 시뮬레이션", .75, 3.25, 6.5, .85, 17, "C6D7D1")
    add_rect(slide, .75, 5.35, 6.6, .02, LIME, LIME)
    add_text(slide, "기업 의사결정 관점 + 분석·구현 근거", .75, 5.58, 6, .3, 11, LIME, True)
    add_text(slide, "2026. 08. 28", .75, 6.48, 2.2, .25, 10, "9DB7AE")

    # 2. Executive summary
    slide = add_slide_base(prs, "Executive summary", 2, "예측에서 비교까지, 운행 계획 검토의 근거를 만들었습니다",
                           "저장된 최종 모델을 웹서비스에 연결해 재학습 없이 즉시 추론합니다.")
    cards = [
        ("01", "문제", "운행 전 에너지 소비량을\n정량적으로 가늠하기 어려움"),
        ("02", "해법", "9개 운행 조건 기반\nXGBoost 회귀 예측"),
        ("03", "서비스", "현재·대안 조건 비교와\n예상 절감률 시뮬레이션"),
        ("04", "성과", "테스트 R² 0.948\n평균 오차 0.660 kWh/100km"),
    ]
    for i, (num, label, body) in enumerate(cards):
        x = .68 + i*3.12
        add_rect(slide, x, 2.28, 2.82, 2.5, WHITE, LINE, True)
        add_text(slide, num, x+.22, 2.52, .5, .3, 11, GREEN, True)
        add_text(slide, label, x+.22, 2.98, 2.3, .35, 18, INK, True)
        add_text(slide, body, x+.22, 3.55, 2.34, .75, 12, MUTED)
    add_rect(slide, .68, 5.18, 12, 1.18, PALE, PALE, True)
    add_text(slide, "핵심 판단", .95, 5.47, 1.2, .3, 11, GREEN, True)
    add_text(slide, "본 결과는 조건별 비교를 위한 모델 기반 예상치이며, 실제 운영 적용 전 실차 데이터 검증이 필요합니다.",
             2.1, 5.42, 9.95, .42, 15, INK, True)

    # 3. 문제 정의
    slide = add_slide_base(prs, "01 · Business problem", 3, "운행 계획 수립 시 소비량 불확실성을 줄여야 합니다",
                           "업무 이슈 → 예측 대상 → 분석 목표 → 활용 방안의 흐름으로 문제를 재정의했습니다.")
    issues = [
        ("충전 계획 오차", "주행 조건 차이를 반영하지 못하면 필요한 에너지를 과소·과대 추정"),
        ("운영비 판단 지연", "대안 조건의 소비량 차이를 사전에 비교할 정량 기준 부족"),
        ("설명 근거 부족", "예측 결과와 함께 어떤 조건이 모델에서 중요했는지 전달 필요"),
    ]
    for i,(head,body) in enumerate(issues):
        y=2.25+i*1.18
        add_text(slide, f"0{i+1}", .75, y, .45, .3, 11, GREEN, True)
        add_text(slide, head, 1.35, y-.02, 2.2, .35, 16, INK, True)
        add_text(slide, body, 3.55, y, 4.0, .5, 11, MUTED)
        add_line(slide, .75, y+.76, 7.45, y+.76, LINE)
    add_rect(slide, 8.05, 2.22, 4.62, 3.7, DARK, DARK, True)
    add_text(slide, "분석 목표", 8.42, 2.62, 2, .3, 11, LIME, True)
    add_text(slide, "주행·환경 조건으로\n에너지 소비량을 예측하고\n계획 간 비교의 객관성을 높인다", 8.4, 3.18, 3.85, 1.55, 22, WHITE, True)
    add_text(slide, "사용자  · 차량 관제 / 충전 운영 / 운영 기획", 8.4, 5.25, 3.7, .3, 9, "B7CBC4")

    # 4. 분석 흐름
    slide = add_slide_base(prs, "01 · Project flow", 4, "분석 결과를 실제 사용 가능한 추론 흐름으로 연결했습니다")
    steps = [
        ("01", "DATA", "8,000건\n10개 변수"),
        ("02", "EDA", "품질·분포\n상관관계"),
        ("03", "MODEL", "7개 후보\nCV 비교"),
        ("04", "SERVE", "joblib\nFastAPI"),
        ("05", "DECIDE", "현재·대안\n조건 비교"),
    ]
    for i,(num,label,body) in enumerate(steps):
        x=.68+i*2.52
        add_rect(slide,x,2.35,2.12,2.55,DARK if i==4 else WHITE,DARK if i==4 else LINE,True)
        add_text(slide,num,x+.2,2.6,.5,.25,10,LIME if i==4 else GREEN,True)
        add_text(slide,label,x+.2,3.03,1.7,.3,14,WHITE if i==4 else INK,True)
        add_text(slide,body,x+.2,3.62,1.7,.7,13,"C7D8D2" if i==4 else MUTED)
        if i<4:
            add_text(slide,"→",x+2.15,3.4,.35,.35,17,GREEN,True,PP_ALIGN.CENTER)
    add_text(slide, "운영 원칙", .72, 5.5, 1.2, .25, 10, GREEN, True)
    add_text(slide, "웹 요청마다 학습하지 않고, 검증된 최종 모델을 시작 시 한 번 로드해 재사용", 1.8, 5.42, 10.2, .42, 15, INK, True)

    # 5. 데이터 구조
    slide = add_slide_base(prs, "02 · Data", 5, "8,000건의 주행·환경 조건으로 연속형 소비량을 예측합니다")
    add_metric(slide,.7,2.2,2.65,1.65,"관측 데이터","8,000","건",True,inline_unit=True)
    add_metric(slide,3.55,2.2,2.65,1.65,"입력 Feature","9","개",inline_unit=True)
    add_metric(slide,6.4,2.2,2.65,1.65,"결측값","0","건",inline_unit=True)
    add_metric(slide,9.25,2.2,3.35,1.65,"예측 Target","kWh/100km","")
    groups = [
        ("주행", "평균 속도 · 운전 성향 · 주행 거리"),
        ("차량", "적재량 · 배터리 온도 · 타이어 공기압"),
        ("환경", "외기 온도 · 공조 전력 · 도로 경사도"),
    ]
    for i,(label,body) in enumerate(groups):
        x=.7+i*4.02
        add_rect(slide,x,4.45,3.76,1.35,WHITE,LINE,True)
        add_text(slide,label,x+.22,4.69,.75,.28,11,GREEN,True)
        add_text(slide,body,x+.22,5.13,3.3,.35,11,INK,True)
    add_text(slide,"문제 유형  ·  지도학습 기반 회귀",.72,6.25,4,.25,10,MUTED)

    # 6. EDA
    slide = add_slide_base(prs, "02 · EDA", 6, "데이터 품질은 안정적이며, 소비량은 넓은 연속 구간에 분포합니다",
                           "결측 0건 · 중복 0건 · 입력 변수별 관측 범위 확인")
    add_image_contain(slide,ASSET_DIR/"target_distribution.png",.68,2.17,7.25,4.35)
    add_rect(slide,8.25,2.17,4.4,4.35,WHITE,LINE,True)
    add_text(slide,"해석",8.58,2.52,1,.3,11,GREEN,True)
    add_text(slide,"타깃 값이 특정 한 점에 몰리지 않아\n회귀 모델이 다양한 소비량 구간을\n학습할 수 있습니다.",8.57,3.04,3.55,1.0,17,INK,True)
    add_line(slide,8.57,4.38,12.18,4.38,LINE)
    add_text(slide,"관측 범위 밖 입력",8.58,4.7,2,.25,10,GOLD,True)
    add_text(slide,"예측은 허용하되 외삽 불확실성\n경고를 화면에 명시합니다.",8.57,5.12,3.4,.7,12,MUTED)

    # 7. correlation
    slide = add_slide_base(prs, "02 · EDA", 7, "경사도·적재량·속도가 소비량과 높은 양의 상관을 보입니다",
                           "상관관계는 탐색적 관련성이며 인과관계를 의미하지 않습니다.")
    add_image_fit(slide,ASSET_DIR/"correlations.png",.68,2.15,7.6,4.45)
    insights=[
        ("도로 경사도", "+0.50", "상관계수 최상위"),
        ("적재량", "+0.47", "차량 부하 관련"),
        ("평균 속도", "+0.45", "주행 조건 관련"),
    ]
    for i,(label,value,note) in enumerate(insights):
        y=2.2+i*1.22
        add_rect(slide,8.62,y,3.98,1.0,WHITE,LINE,True)
        add_text(slide,label,8.88,y+.18,1.65,.25,12,INK,True)
        add_text(slide,value,10.72,y+.15,1.1,.35,18,GREEN,True,PP_ALIGN.RIGHT)
        add_text(slide,note,8.88,y+.57,2.8,.2,8,MUTED)
    add_rect(slide,8.62,5.98,3.98,.48,PALE,PALE,True)
    add_text(slide,"기온·공기압 등은 약한 음의 상관",8.85,6.11,3.5,.2,9,GREEN,True)

    # 8. Model compare
    slide = add_slide_base(prs, "03 · Modeling", 8, "비교 기준을 세우고 튜닝 XGBoost를 최종 후보로 선정했습니다",
                           "Train 80% / Test 20% · 학습 데이터 내 5-Fold 교차검증 · RMSE 기준 튜닝")
    add_image_contain(slide,ASSET_DIR/"model_comparison.png",.68,2.15,8.1,4.4)
    add_rect(slide,9.02,2.15,3.63,4.4,DARK,DARK,True)
    add_text(slide,"선정 근거",9.35,2.5,1.5,.3,11,LIME,True)
    add_text(slide,"XGBoost\nTuned",9.34,3.05,2.5,.72,24,WHITE,True)
    add_text(slide,"CV RMSE",9.35,4.05,1.4,.25,10,"B8CCC5")
    add_text(slide,"0.887",9.35,4.45,2.2,.55,28,LIME,True)
    add_text(slide,"비선형 관계·변수 간 상호작용을\n반영하며 후보 중 오차가 가장 낮음",9.35,5.35,2.75,.7,11,"C8D8D3")

    # 9. performance
    slide = add_slide_base(prs, "03 · Performance", 9, "최종 모델은 테스트 데이터 변동의 94.8%를 설명했습니다")
    add_image_fit(slide,ASSET_DIR/"actual_vs_predicted.png",.68,2.1,6.7,4.55)
    add_metric(slide,7.73,2.18,2.22,1.62,"R²","0.948","설명력",True)
    add_metric(slide,10.18,2.18,2.35,1.62,"MAE","0.660","kWh/100km")
    add_metric(slide,7.73,4.08,2.22,1.62,"RMSE","0.834","kWh/100km")
    add_rect(slide,10.18,4.08,2.35,1.62,PALE,PALE,True)
    add_text(slide,"해석",10.42,4.33,1,.25,10,GREEN,True)
    add_text(slide,"평균 절대 오차가\n1 kWh/100km 미만",10.4,4.82,1.8,.56,13,INK,True)
    add_text(slide,"※ 교육용 데이터 기준 성능",7.75,6.17,4.2,.2,9,MUTED)

    # 10 importance
    slide = add_slide_base(prs, "03 · Interpretation", 10, "도로 경사도·적재량·속도가 모델 예측의 상위 요인입니다",
                           "Feature Importance는 예측 기여의 상대적 비중이며 방향성과 인과를 설명하지 않습니다.")
    add_image_fit(slide,ASSET_DIR/"feature_importance.png",.68,2.14,7.6,4.45)
    add_rect(slide,8.6,2.15,4.03,2.1,DARK,DARK,True)
    add_text(slide,"표현 원칙",8.92,2.5,1.2,.25,10,LIME,True)
    add_text(slide,"“높은 평균 속도는 높은 소비량\n예측과 관련된 주요 조건 중 하나”",8.9,3.0,3.3,.75,15,WHITE,True)
    add_rect(slide,8.6,4.52,4.03,1.55,"FFF7DF","E7C46D",True)
    add_text(slide,"주의",8.9,4.82,.7,.25,10,GOLD,True)
    add_text(slide,"“속도 때문에 소비량이 증가”와 같은\n인과 표현은 사용하지 않습니다.",8.9,5.2,3.25,.55,11,INK)

    # 11 architecture
    slide = add_slide_base(prs, "04 · Web service", 11, "모델 추론과 웹 요청 처리를 분리해 재사용성과 안정성을 확보했습니다")
    layers=[
        ("USER", "운행 조건 입력\n현재 / 대안 계획"),
        ("FASTAPI", "Route · Schema\n범위 이탈 경고"),
        ("SERVICE", "예측 · 총에너지\n절감률 계산"),
        ("MODEL", "joblib 1회 로드\nXGBoost Pipeline"),
    ]
    for i,(label,body) in enumerate(layers):
        x=.72+i*3.05
        add_rect(slide,x,2.35,2.55,2.3,DARK if i==3 else WHITE,DARK if i==3 else LINE,True)
        add_text(slide,label,x+.23,2.68,1.9,.3,12,LIME if i==3 else GREEN,True)
        add_text(slide,body,x+.23,3.31,2.0,.7,14,WHITE if i==3 else INK,True)
        if i<3: add_text(slide,"→",x+2.57,3.25,.45,.4,18,GREEN,True,PP_ALIGN.CENTER)
    principles=["저장된 feature_columns 순서 유지","pathlib + 기존 MODEL_DIR 재사용","요청마다 모델 재학습·재로드 금지"]
    for i,item in enumerate(principles):
        x=.72+i*4.05
        add_rect(slide,x,5.18,3.72,.82,PALE,PALE,True)
        add_text(slide,f"0{i+1}",x+.2,5.43,.35,.2,9,GREEN,True)
        add_text(slide,item,x+.63,5.37,2.83,.32,10,INK,True)

    # 12 UI prediction
    slide = add_slide_base(prs, "04 · Service demo", 12, "9개 운행 조건으로 예상 소비량과 총소비 에너지를 즉시 확인합니다")
    add_image_fit(slide,ASSET_DIR/"web_prediction.png",.68,2.12,8.55,4.55)
    add_rect(slide,9.54,2.13,3.1,4.54,DARK,DARK,True)
    add_text(slide,"입력 예시",9.86,2.5,1.4,.25,10,LIME,True)
    add_text(slide,"속도 70 km/h\n적재량 180 kg\n주행 거리 50 km",9.85,2.98,2.2,1.0,14,WHITE,True)
    add_line(slide,9.86,4.25,12.24,4.25,"557A70")
    add_text(slide,"예상 소비량",9.86,4.57,1.5,.25,10,"B8CCC5")
    add_text(slide,"20.34",9.85,4.98,1.8,.5,27,LIME,True)
    add_text(slide,"kWh/100km",9.87,5.48,1.5,.22,9,"C8D8D3")
    add_text(slide,"총 10.17 kWh",9.86,5.98,1.8,.25,11,WHITE,True)

    # 13 comparison warning
    slide = add_slide_base(prs, "04 · Service demo", 13, "현재·대안 계획을 비교하고 외삽 입력에는 불확실성 경고를 제공합니다",
                           "범위 밖 입력도 시뮬레이션하되 결과 해석 위험을 사용자에게 명시합니다.")
    add_image_fit(slide,ASSET_DIR/"web_simulation_warning.png",.68,2.12,7.45,4.6)
    add_rect(slide,8.47,2.13,4.16,1.25,WHITE,LINE,True)
    add_text(slide,"조건별 비교",8.77,2.43,1.3,.25,11,GREEN,True)
    add_text(slide,"소비량 · 총에너지 · 차이 · 예상 절감률",8.77,2.84,3.2,.28,11,INK,True)
    add_rect(slide,8.47,3.67,4.16,1.55,"FFF7DF","E7C46D",True)
    add_text(slide,"경고 정책",8.77,3.97,1.2,.25,11,GOLD,True)
    add_text(slide,"관측 범위 밖 항목과 범위를 구체적으로 표시\n→ 모델 외삽 결과임을 인지한 상태에서 판단",8.77,4.4,3.4,.55,11,INK)
    add_rect(slide,8.47,5.51,4.16,1.18,DARK,DARK,True)
    add_text(slide,"절감률은 대안 소비량이 더 낮을 때만 계산",8.77,5.9,3.35,.35,11,WHITE,True)

    # 14 limits
    slide = add_slide_base(prs, "05 · Risk & limits", 14, "높은 모델 성능과 실제 운영 효과는 동일하지 않습니다",
                           "실습용 데이터와 모델 기반 예측이라는 경계를 분명히 관리해야 합니다.")
    risks=[
        ("데이터 대표성", "교육용 8,000건이 실제 차종·노화·기상 조건을 모두 대표하지 않음", "실차·차종별 데이터 수집"),
        ("외삽 불확실성", "학습 관측 범위 밖에서는 모델의 오차 수준을 보장할 수 없음", "경고 + 운영 적용 범위 정의"),
        ("비인과적 해석", "Feature Importance만으로 변수 변화의 원인·효과를 단정할 수 없음", "SHAP·실험 설계 추가 검토"),
        ("운영 검증 부재", "예상 절감률은 시뮬레이션이며 실제 절감 성과가 아님", "파일럿 A/B 검증 및 모니터링"),
    ]
    headers=[("리스크",.72,2.23,2.15),("현재 판단",2.9,2.23,5.75),("대응 방향",8.92,2.23,3.7)]
    for text,x,y,w in headers: add_text(slide,text,x,y,w,.3,11,GREEN,True)
    for i,(risk,current,response) in enumerate(risks):
        y=2.72+i*.91
        add_rect(slide,.7,y,11.95,.76,WHITE,LINE,True)
        add_text(slide,risk,.92,y+.21,1.7,.25,11,INK,True)
        add_text(slide,current,2.92,y+.17,5.5,.38,10,MUTED)
        add_text(slide,response,8.95,y+.2,3.25,.28,10,GREEN,True)
    add_text(slide,"운영 메시지  ·  예상 소비량 / 모델 기반 예측 / 조건별 시뮬레이션 / 예상 절감률",.73,6.55,9,.25,10,MUTED)

    # 15 conclusion
    slide = add_slide_base(prs, "05 · Conclusion", 15, "다음 단계는 실차 데이터로 ‘예측 가능성’을 ‘업무 효과’로 검증하는 것입니다", dark=True)
    roadmap=[
        ("NOW", "실습 모델·웹서비스", "저장 모델 추론\n조건 비교·경고"),
        ("NEXT", "파일럿 검증", "실차 오차 측정\n사용자 피드백"),
        ("SCALE", "운영 확장", "차종·배터리 상태\n충전·관제 연동"),
    ]
    for i,(tag,title,body) in enumerate(roadmap):
        x=.72+i*4.12
        add_rect(slide,x,2.25,3.65,2.33,"174E41","2B6657",True)
        add_text(slide,tag,x+.28,2.58,1,.25,10,LIME,True)
        add_text(slide,title,x+.28,3.08,2.8,.35,17,WHITE,True)
        add_text(slide,body,x+.28,3.67,2.8,.62,12,"C4D5CF")
        if i<2: add_text(slide,"→",x+3.68,3.22,.4,.35,17,LIME,True,PP_ALIGN.CENTER)
    add_rect(slide,.72,5.15,11.9,1.05,LIME,LIME,True)
    add_text(slide,"제언",1.02,5.47,.8,.25,11,DARK,True)
    add_text(slide,"의사결정 지원 도구로 제한된 파일럿을 운영하고, 실제 오차와 절감 효과를 함께 측정하십시오.",
             1.82,5.4,10.1,.4,15,DARK,True)

    prs.core_properties.title = "전기차 에너지 소비량 예측 웹서비스 프로젝트 보고서"
    prs.core_properties.subject = "머신러닝 분석 및 FastAPI 웹서비스 구현"
    prs.core_properties.author = "energy_consumption project"
    prs.core_properties.comments = "교육 및 실습 목적의 모델 기반 예측 보고서"
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    output = build()
    print(output)
